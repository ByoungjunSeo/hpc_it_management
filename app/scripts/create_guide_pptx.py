#!/usr/bin/env python3
"""Create IT Asset Management User Guide PPTX with screenshots."""
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

SCREENSHOTS = os.path.join(os.path.dirname(__file__), '..', 'data', 'screenshots')
OUTPUT = os.path.join(os.path.dirname(__file__), '..', 'data', 'IT_자산관리_사용가이드.pptx')

# Color scheme
PRIMARY = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Layout constants for content slides
HEADER_H = Inches(0.9)
IMG_TOP = Inches(1.05)
IMG_LEFT = Inches(0.3)
IMG_BOT_MARGIN = Inches(0.2)
IMG_MAX_H = prs.slide_height - IMG_TOP - IMG_BOT_MARGIN  # ~6.25"
IMG_MAX_W = Inches(6.3)  # about half the slide
BULLET_LEFT = Inches(6.9)
BULLET_TOP = Inches(1.15)
BULLET_WIDTH = Inches(6.1)
BULLET_HEIGHT = Inches(5.8)
# Target aspect ratio for image area (h/w)
TARGET_RATIO = IMG_MAX_H / IMG_MAX_W  # ~0.99


def add_bg_rect(slide, color=LIGHT_BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Malgun Gothic'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_screenshot(slide, filename, left, top, max_width, max_height):
    """Add screenshot, cropping tall images from top to fit the target area."""
    path = os.path.join(SCREENSHOTS, filename)
    if not os.path.exists(path):
        print(f"  WARNING: {filename} not found, skipping")
        return None

    img = Image.open(path)
    w, h = img.size
    src_ratio = h / w  # source aspect ratio
    tgt_ratio = max_height / max_width  # target area ratio (in EMU)

    # Convert EMU ratio to pixel-comparable ratio
    tgt_ratio_px = (max_height / max_width)  # Both are EMU so ratio is unitless

    use_path = path
    tmp_file = None

    if src_ratio > tgt_ratio_px * 1.15:
        # Image is significantly taller than target area -> crop from top
        crop_h = int(w * tgt_ratio_px)
        if crop_h < h:
            cropped = img.crop((0, 0, w, crop_h))
            tmp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            cropped.save(tmp_file.name)
            use_path = tmp_file.name
            w, h = w, crop_h
            src_ratio = h / w
            print(f"  Cropped {filename}: {img.size[1]}px -> {crop_h}px height")

    # Scale to fit within max_width x max_height
    scale_w = max_width / w
    scale_h = max_height / h
    scale = min(scale_w, scale_h)
    img_width = int(w * scale)
    img_height = int(h * scale)

    pic = slide.shapes.add_picture(use_path, left, top, img_width, img_height)

    if tmp_file:
        try:
            os.unlink(tmp_file.name)
        except OSError:
            pass

    return pic


# ══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, PRIMARY)
add_rect(slide, 0, Inches(5.2), prs.slide_width, Inches(2.3), RGBColor(0x1D, 0x4E, 0xD8))

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "IT 자산 관리 시스템", font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "사용자 가이드", font_size=28, color=RGBColor(0xBF, 0xDB, 0xFE),
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(1.2),
             "서버실 장비 현황 ・ 입출고 관리 ・ 부품 추적 ・ IP 관리\n네트워크 레이아웃 ・ 전력 분전반 ・ GPU 모니터링 ・ 이력 관리",
             font_size=16, color=RGBColor(0x93, 0xC5, 0xFD), alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 2: Table of Contents
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide)
add_rect(slide, 0, 0, prs.slide_width, Inches(1.0), PRIMARY)
add_text_box(slide, Inches(0.5), Inches(0.2), Inches(5), Inches(0.6),
             "목차", font_size=28, color=WHITE, bold=True)

toc_items = [
    ("1", "로그인", "시스템 접속 및 권한 관리"),
    ("2", "대시보드", "전체 자산 현황 한눈에 보기"),
    ("3", "입출고 관리", "입고 등록 → 사용 등록 → 반납 프로세스"),
    ("4", "자산 현황", "전체 자산 목록 및 상세 정보"),
    ("5", "부품 현황", "CPU/메모리/GPU 등 부품 재고 관리"),
    ("6", "서버실 관리", "서버실 → 랙 → 장비 위치 관리 + 전원 상태"),
    ("7", "사무실 / 장비실", "사무실 및 보관 장비 관리"),
    ("8", "전력 분전반", "분전반 → PDU → 장비 전력 트리"),
    ("9", "네트워크 레이아웃", "스위치-장비 포트 연결 관리"),
    ("10", "IP 관리", "IP 주소 대역별 할당 현황"),
    ("11", "자산 매칭", "SSH 스캔으로 HW 자동 탐색/비교"),
    ("12", "GPU 모니터링", "실시간 GPU 온도/사용률/전력 모니터링"),
    ("13", "이력 관리", "전체 작업 이력 감사 로그"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.3) + Inches(i * 0.43)
    add_rect(slide, Inches(0.8), y, Inches(0.5), Inches(0.35), PRIMARY)
    add_text_box(slide, Inches(0.8), y + Pt(2), Inches(0.5), Inches(0.35),
                 num, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), y + Pt(1), Inches(3), Inches(0.35),
                 title, font_size=15, color=DARK, bold=True)
    add_text_box(slide, Inches(4.5), y + Pt(2), Inches(7), Inches(0.35),
                 desc, font_size=13, color=GRAY)


# ══════════════════════════════════════════════════════════
# Content slides with screenshots
# ══════════════════════════════════════════════════════════

slides_data = [
    {
        "title": "1. 로그인",
        "screenshot": "01_login.png",
        "bullets": [
            "아이디와 비밀번호를 입력하여 로그인",
            "권한 등급: 관리자 / 유지보수 / 조회",
            "관리자: 모든 기능 사용 가능",
            "유지보수: 등록/수정/삭제 가능",
            "조회: 읽기 전용 (변경 불가)",
        ]
    },
    {
        "title": "2. 대시보드",
        "screenshot": "02_dashboard.png",
        "bullets": [
            "로그인 후 처음 보이는 메인 화면",
            "자산 유형별 현황 (서버/스위치/스토리지 등)",
            "서버실별 랙 사용률 요약",
            "최근 입출고/변경 이력",
            "통합 검색: 장비명, IP, 사용자 등",
        ]
    },
    # ---- 입출고 관리 섹션 ----
    {
        "title": "3. 입출고 관리 — 목록",
        "screenshot": "step_in_01_list.png",
        "bullets": [
            "장비 및 부품의 입고/출고를 기록 관리",
            "입출고 이력 조회 및 필터링",
            "",
            "[핵심 워크플로우]",
            "  신규 장비 → 입고 등록 → 사용 등록",
            "  기존 자산 → 사용 등록 (바로)",
            "",
            "[화면 안내]",
            "  ① 녹색 '+ 입고 등록' 버튼 클릭",
            "     → 신규 장비를 시스템에 최초 등록",
        ]
    },
    {
        "title": "3-1. [입고 등록] STEP 1 — 자산유형 및 소유 선택",
        "screenshot": "step_in_02_type.png",
        "step": True,
        "bullets": [
            "입고 등록 폼 상단에서 기본 분류 선택",
            "",
            "[화면 안내]",
            "  ① 자산유형 선택",
            "     - 장비: 서버, 스위치, 스토리지 등",
            "     - 부품: CPU, 메모리, 디스크, PSU 등",
            "",
            "  ② 소유구분 선택",
            "     - 자사: 직접 구매한 장비",
            "     - 업체: 외부 업체에서 반입된 장비",
            "",
            "※ 모든 신규 장비(구매/업체 불문)는",
            "   반드시 입고 등록을 먼저 해야 합니다.",
        ]
    },
    {
        "title": "3-2. [입고 등록] STEP 2 — 기본 정보 입력",
        "screenshot": "step_in_03_fill.png",
        "step": True,
        "bullets": [
            "자산유형/소유 선택 후 상세 정보 입력",
            "",
            "[화면 안내]",
            "  ① 관리번호 (필수)",
            "     - 체계에 맞게 부여 (예: TPC-SV-2U-99)",
            "     - 추천 번호가 자동 표시됩니다",
            "",
            "  ② 자산번호 (선택)",
            "  ③ 제조사",
            "  ④ 모델명",
            "  ⑤ 시리얼번호",
            "  ⑥ U 사이즈 (서버/스토리지/스위치)",
            "",
            "  ⑦ '입고 등록' 버튼 클릭하여 완료",
        ]
    },
    {
        "title": "3-3. [입고 등록] STEP 3 — 업체 장비 등록",
        "screenshot": "step_in_04_vendor.png",
        "step": True,
        "bullets": [
            "소유구분 '업체' 선택 시 추가 입력",
            "",
            "[화면 안내]",
            "  ① 소유구분: '업체' 선택",
            "",
            "  ② 업체 선택 (드롭다운 표시됨)",
            "     - 기존 등록 업체에서 선택",
            "     - 또는 '+ 신규 업체 등록'으로 추가",
            "",
            "[업체 장비 특이사항]",
            "  - 관리번호가 업체명 기반으로 자동 생성",
            "  - 예: 클라우드원-001, 엘텍코리아-003",
            "  - 반납예정일 입력 가능",
        ]
    },
    {
        "title": "3-4. [사용 등록] STEP 1 — 기존 자산 로딩",
        "screenshot": "step_use_02_mgmt.png",
        "step": True,
        "bullets": [
            "입출고 목록에서 파란색 '+ 사용 등록' 클릭",
            "",
            "[화면 안내]",
            "  ① 관리번호 입력",
            "     - 타이핑 시 자동완성 목록 표시",
            "     - 선택 시 기존 정보 자동 로딩",
            "     (모델명, IP, 접속정보, 하드웨어 등)",
            "",
            "[사용 시나리오]",
            "  - 신규 입고된 장비의 첫 사용 등록",
            "  - 기존 장비의 위치 변경/정보 업데이트",
            "",
            "※ 신규 장비는 반드시 입고 등록 먼저!",
        ]
    },
    {
        "title": "3-5. [사용 등록] STEP 2 — 하드웨어 정보",
        "screenshot": "step_use_03_hardware.png",
        "step": True,
        "bullets": [
            "장비에 장착된 부품 정보를 입력",
            "",
            "[화면 안내]",
            "  ① 부품 추가 버튼",
            "     + CPU / + 메모리 / + 디스크 등 클릭",
            "",
            "  ② 부품 행 입력",
            "     - 소유: 자사 / 업체",
            "     - 부품코드: 드롭다운 또는 직접 입력",
            "     - 수량: 장착 수량 입력",
            "",
            "[PSU (전원장치)]",
            "  - '+ PSU' 버튼으로 추가",
            "  - Active / Standby 역할 선택",
        ]
    },
    {
        "title": "3-6. [사용 등록] STEP 3 — IP / 접속정보",
        "screenshot": "step_use_04_ip_cred.png",
        "step": True,
        "bullets": [
            "네트워크 및 접속 정보를 입력",
            "",
            "[화면 안내]",
            "  ① 접속 정보 (계정)",
            "     - Root / BMC / OS / 기타 계정 등록",
            "     - 사용자명, 비밀번호 입력",
            "",
            "  ② IP 주소",
            "     - 용도: Management / BMC / IB 등",
            "     - 인터페이스 유형 및 속도 선택",
            "     - IP 주소 입력 (예: 10.100.230.50)",
        ]
    },
    {
        "title": "3-7. [사용 등록] STEP 4 — 위치 배정",
        "screenshot": "step_use_05_location.png",
        "step": True,
        "bullets": [
            "장비를 배치할 서버실/랙 위치 지정",
            "",
            "[화면 안내]",
            "  ① 서버실 선택 (드롭다운)",
            "",
            "  ② 랙 선택 (서버실에 따라 필터됨)",
            "",
            "  ③ 랙 미리보기",
            "     - 현재 랙의 장비 배치 현황 표시",
            "     - 빈 위치를 클릭하면 자동 선택",
            "     - U 시작위치 + 크기 자동 설정",
            "",
            "→ 모든 입력 후 하단 '등록' 버튼 클릭",
        ]
    },
    {
        "title": "3-8. 장비 이력 상세",
        "screenshot": "03e_equip_detail.png",
        "bullets": [
            "입출고 목록에서 관리번호 클릭",
            "",
            "[확인 가능한 정보]",
            "  - 기본 정보 (관리번호, 모델명 등)",
            "  - 하드웨어 구성 (CPU, 메모리 등)",
            "  - IP 주소 목록",
            "  - 접속 정보 (계정)",
            "  - 입출고 이력 (시간순)",
            "",
            "[상태 흐름]",
            "  입고 → 사용중 → 반납완료",
            "  (필요 시 재입고/재사용 가능)",
        ]
    },
    # ---- 신청서 ----
    {
        "title": "4. 신청서 관리",
        "screenshot": "04_requests.png",
        "bullets": [
            "장비 반출/반입/변경 등의 신청서 관리",
            "신청 → 승인 → 처리 워크플로우",
            "신청서 상태별 필터링 및 조회",
        ]
    },
    # ---- 자산 현황 ----
    {
        "title": "5. 자산 현황",
        "screenshot": "05_assets.png",
        "bullets": [
            "전체 자산 목록을 한눈에 조회",
            "서버실/유형/상태별 필터링",
            "검색: 관리번호, 모델명, IP, 사용자",
            "자산 클릭 → 상세 정보 페이지 이동",
            "상세: IP 목록, 인증정보, 컴퓨팅 모듈 확인",
        ]
    },
    {
        "title": "5-1. 자산 상세",
        "screenshot": "16_asset_detail.png",
        "bullets": [
            "자산의 모든 상세 정보를 확인/수정",
            "기본 정보: 모델, 시리얼, 위치, 상태",
            "IP 주소: 관리IP, BMC, OS, IB 등 다중 등록",
            "인증 정보: Root, User, BMC 계정 관리",
            "컴퓨팅 모듈: CPU, 메모리, 디스크, GPU, PSU",
            "대여 현황 및 이력 확인",
        ]
    },
    # ---- 부품 현황 ----
    {
        "title": "6. 부품 현황",
        "screenshot": "06_modules.png",
        "bullets": [
            "CPU, 메모리, GPU, 디스크 등 부품별 재고 현황",
            "설치/가용/전체 수량 한눈에 확인",
            "부품 코드별 상세 이력 추적",
            "자산 매칭(스캔)과 연동하여 자동 갱신",
        ]
    },
    # ---- 서버실 관리 ----
    {
        "title": "7. 서버실 관리",
        "screenshot": "07_rooms.png",
        "bullets": [
            "서버실 목록 및 각 서버실의 랙 구성 조회",
            "랙 내 장비 유형별 색상 구분",
            "전원 확인: BMC를 통한 전원 ON/OFF 상태 표시",
            "전원 상태: 초록(ON), 빨강(OFF), 검정(BMC 미등록)",
            "U 단위 위치 기반 장비 시각화",
            "랙 드래그앤드롭으로 순서 변경 가능",
        ]
    },
    {
        "title": "8. 사무실 / 장비실",
        "screenshot": "08_offices.png",
        "bullets": [
            "사무실: PC, 모니터 등 사무용 장비 관리",
            "장비실: 보관 중인 여분 장비 관리",
            "위치별, 사용자별 조회 가능",
        ]
    },
    {
        "title": "9. 전력 분전반",
        "screenshot": "10_power_panel.png",
        "bullets": [
            "전력 공급 경로를 트리 구조로 시각화",
            "메인 분전반 → 분전반 → UPS → PDU 계층",
            "각 노드에 연결된 장비/랙 확인",
            "노드별 용량, 상태, 메모 관리",
        ]
    },
    {
        "title": "10. 네트워크 레이아웃",
        "screenshot": "11_network.png",
        "bullets": [
            "스위치-장비 간 물리적 포트 연결 관리",
            "서버실별 네트워크 구성도 시각화",
            "포트별 케이블 종류, 속도, 상태 기록",
            "연결 추가/수정/삭제",
        ]
    },
    {
        "title": "11. IP 관리",
        "screenshot": "12_ip_management.png",
        "bullets": [
            "서브넷 대역별 IP 할당 현황 시각화",
            "사용 중 / 미사용 IP 한눈에 확인",
            "IP 클릭 → 해당 자산으로 이동",
            "AIDC, HPC, Office 등 존(Zone)별 필터",
        ]
    },
    {
        "title": "12. 자산 매칭 (SSH 스캔)",
        "screenshot": "13_discovery.png",
        "bullets": [
            "서버에 SSH 접속하여 HW 정보 자동 탐색",
            "CPU, 메모리, 디스크, GPU, 네트워크 카드 탐지",
            "등록된 정보와 실제 HW 비교 (일치/불일치/누락)",
            "스캔 결과를 자산에 바로 반영 가능",
            "개별 또는 전체 일괄 스캔 지원",
        ]
    },
    {
        "title": "13. GPU 모니터링",
        "screenshot": "14_gpu_monitoring.png",
        "bullets": [
            "GPU 보유 서버의 실시간 상태 모니터링",
            "도넛 차트: 평균 온도/사용률/메모리/전력",
            "총 전력 소비량 (kW) 표시",
            "서버별 GPU 게이지 바 (온도/사용률/메모리/전력)",
            "dcgmi 우선, nvidia-smi 폴백 자동 전환",
            "자동 갱신 (30초 간격) 지원",
            "서버실별 필터링 가능",
        ]
    },
    {
        "title": "14. 이력 관리",
        "screenshot": "15_audit_log.png",
        "bullets": [
            "모든 등록/수정/삭제/스캔 작업의 감사 로그",
            "작업자, 일시, 대상 자산, 변경 내용 기록",
            "날짜 범위 및 작업 유형별 필터링",
        ]
    },
]

STEP_ACCENT = RGBColor(0x3B, 0x82, 0xF6)

for sd in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide)

    is_step = sd.get("step", False)

    # Header bar
    add_rect(slide, 0, 0, prs.slide_width, HEADER_H, PRIMARY)

    if is_step:
        # Step badge + title
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.17), Inches(1.0), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = STEP_ACCENT
        badge.line.fill.background()
        tf_b = badge.text_frame
        p_b = tf_b.paragraphs[0]
        # Extract step part from title (e.g., "STEP 1")
        title_text = sd["title"]
        p_b.text = title_text.split("—")[0].split("STEP")[0].strip().split(".")[-1].strip() if "STEP" in title_text else ""
        step_part = title_text.split("STEP")
        if len(step_part) > 1:
            step_num = step_part[1].split("—")[0].strip()
            p_b.text = "STEP " + step_num
        p_b.font.size = Pt(13)
        p_b.font.bold = True
        p_b.font.color.rgb = WHITE
        p_b.alignment = PP_ALIGN.CENTER
        # Title after badge
        title_after = "—" + title_text.split("—", 1)[1] if "—" in title_text else title_text
        add_text_box(slide, Inches(1.5), Inches(0.18), Inches(11), Inches(0.6),
                     title_after.strip("— "), font_size=22, color=WHITE, bold=True)
    else:
        add_text_box(slide, Inches(0.5), Inches(0.18), Inches(10), Inches(0.6),
                     sd["title"], font_size=26, color=WHITE, bold=True)

    # Screenshot on the left (~50% of slide), fill height
    pic = add_screenshot(slide, sd["screenshot"], IMG_LEFT, IMG_TOP, IMG_MAX_W, IMG_MAX_H)

    # Add a subtle border around screenshot
    if pic:
        pic.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        pic.line.width = Pt(1)

    # Bullet points on the right
    txBox = slide.shapes.add_textbox(BULLET_LEFT, BULLET_TOP, BULLET_WIDTH, BULLET_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True

    for j, bullet in enumerate(sd["bullets"]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if not bullet:
            # Empty line spacer
            p.text = ""
            p.font.size = Pt(6)
            p.space_after = Pt(2)
        elif bullet.startswith("[") and bullet.endswith("]"):
            # Section header
            p.text = bullet
            p.font.size = Pt(13)
            p.font.color.rgb = STEP_ACCENT
            p.font.name = 'Malgun Gothic'
            p.font.bold = True
            p.space_before = Pt(6)
            p.space_after = Pt(2)
        elif bullet.startswith("※") or bullet.startswith("→"):
            # Highlight / warning
            p.text = bullet
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0xEF, 0x44, 0x44) if bullet.startswith("※") else STEP_ACCENT
            p.font.name = 'Malgun Gothic'
            p.font.bold = True
            p.space_after = Pt(4)
        elif bullet.startswith("  "):
            # Indented sub-item
            p.text = bullet
            p.font.size = Pt(12)
            p.font.color.rgb = GRAY
            p.font.name = 'Malgun Gothic'
            p.space_after = Pt(3)
        else:
            p.text = "•  " + bullet
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.font.name = 'Malgun Gothic'
            p.space_after = Pt(4)


# ══════════════════════════════════════════════════════════
# SLIDE: Tips
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide)
add_rect(slide, 0, 0, prs.slide_width, Inches(0.9), PRIMARY)
add_text_box(slide, Inches(0.5), Inches(0.18), Inches(10), Inches(0.6),
             "활용 팁", font_size=26, color=WHITE, bold=True)

tips = [
    ("통합 검색", "상단 검색창에서 장비명, IP, 사용자명으로 빠른 검색이 가능합니다."),
    ("엑셀 업로드", "입출고 > 엑셀 업로드 기능으로 대량의 자산을 한번에 등록할 수 있습니다."),
    ("자산 매칭", "주기적으로 SSH 스캔을 실행하면 등록 정보와 실제 HW 불일치를 자동 탐지합니다."),
    ("GPU 모니터링", "자동 갱신을 켜두면 30초마다 GPU 상태가 갱신됩니다. 온도 80°C 이상 시 빨간색 경고 표시."),
    ("데이터 백업", "관리자 > 백업 메뉴에서 수동 백업 및 자동 백업(7일 주기) 설정이 가능합니다."),
    ("권한 관리", "관리자는 계정관리에서 사용자 추가 및 권한(관리자/유지보수/조회) 설정이 가능합니다."),
]

for i, (title, desc) in enumerate(tips):
    y = Inches(1.2) + Inches(i * 0.9)
    add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.75), WHITE)
    add_text_box(slide, Inches(0.8), y + Pt(4), Inches(2.5), Inches(0.35),
                 title, font_size=16, color=PRIMARY, bold=True)
    add_text_box(slide, Inches(3.5), y + Pt(6), Inches(8.8), Inches(0.6),
                 desc, font_size=13, color=GRAY)


# ══════════════════════════════════════════════════════════
# SLIDE: End
# ══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, PRIMARY)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
             "감사합니다", font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
             "문의사항은 IT 인프라 관리팀으로 연락해 주세요.", font_size=18,
             color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)


# Save
prs.save(OUTPUT)
print(f"\nPPTX saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
