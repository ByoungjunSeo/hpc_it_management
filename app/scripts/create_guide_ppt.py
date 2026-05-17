#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
IT 자산 관리 시스템 - 상세 사용 가이드 PPT 생성 스크립트
웹사이트를 캡쳐하여 PowerPoint 형태의 가이드를 자동 생성합니다.
"""

import os, sys, time, json
from selenium import webdriver
from selenium.webdriver.firefox.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait, Select
from selenium.webdriver.support import expected_conditions as EC
from PIL import Image as PILImage

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_URL = 'http://localhost:3000'
SCREENSHOT_DIR = '/tmp/guide_screenshots'
OUTPUT_PPT = '/mlcommons_cm/hpc_it_management/IT_자산관리_사용가이드.pptx'

# Colors
BLUE = RGBColor(0x1E, 0x40, 0xAF)
DARK_BLUE = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)

os.makedirs(SCREENSHOT_DIR, exist_ok=True)


def crop_image_if_tall(image_path, max_height_px=900):
    """Crop image to max_height_px if it's taller. Returns (possibly new) path."""
    try:
        img = PILImage.open(image_path)
        w, h = img.size
        if h > max_height_px:
            cropped = img.crop((0, 0, w, max_height_px))
            cropped.save(image_path)
            print('    [CROP] %s: %dx%d -> %dx%d' % (os.path.basename(image_path), w, h, w, max_height_px))
    except Exception as e:
        print('    [CROP WARN] ' + str(e))
    return image_path


def add_fitted_picture(slide, image_path, left, top, max_width, max_height):
    """Add picture to slide, fitting within max_width x max_height while maintaining aspect ratio.
    Vertically centers the image if it doesn't fill the full height."""
    if not image_path or not os.path.exists(image_path):
        return None
    try:
        img = PILImage.open(image_path)
        img_w, img_h = img.size
    except:
        return slide.shapes.add_picture(image_path, left, top, max_width)

    aspect = float(img_w) / float(img_h)

    # Try fitting by width first
    fit_w = max_width
    fit_h = int(fit_w / aspect)

    # If height exceeds max, fit by height instead
    if fit_h > max_height:
        fit_h = max_height
        fit_w = int(fit_h * aspect)

    # Center vertically in available space
    v_offset = (max_height - fit_h) // 2
    actual_top = top + v_offset

    return slide.shapes.add_picture(image_path, left, actual_top, fit_w, fit_h)


BROWSER_WIDTH = 1400
BROWSER_HEIGHT = 1050

# CSS to inject (no longer needed for nav compaction, kept for minor tweaks)
NAV_COMPACT_CSS = ""


def setup_driver():
    geckodriver_path = '/home/mlcommons/.local/lib/python3.6/site-packages/geckodriver_autoinstaller/v0.36.0'
    os.environ['PATH'] = geckodriver_path + ':' + os.environ.get('PATH', '')
    opts = Options()
    opts.headless = True
    opts.add_argument('--width=%d' % BROWSER_WIDTH)
    opts.add_argument('--height=%d' % BROWSER_HEIGHT)
    driver = webdriver.Firefox(options=opts)
    driver.set_window_size(BROWSER_WIDTH, BROWSER_HEIGHT)
    return driver


def inject_nav_css(driver):
    """Inject compact CSS for navigation bar."""
    driver.execute_script("""
        var style = document.getElementById('_guide_nav_css');
        if (!style) {
            style = document.createElement('style');
            style.id = '_guide_nav_css';
            style.textContent = arguments[0];
            document.head.appendChild(style);
        }
    """, NAV_COMPACT_CSS)


def login(driver):
    driver.get(BASE_URL + '/login')
    time.sleep(1)
    driver.find_element(By.NAME, 'username').send_keys('screenshot_temp')
    driver.find_element(By.NAME, 'password').send_keys('screenshot123')
    driver.find_element(By.CSS_SELECTOR, 'button[type="submit"]').click()
    time.sleep(1)


def screenshot(driver, url_or_path, filename, wait=1.5, full_page=False):
    """Take a screenshot and return the file path."""
    if url_or_path and url_or_path.startswith('/'):
        driver.get(BASE_URL + url_or_path)
    time.sleep(wait)
    inject_nav_css(driver)
    time.sleep(0.3)

    filepath = os.path.join(SCREENSHOT_DIR, filename)

    if full_page:
        # Get full page height
        total_height = driver.execute_script("return document.body.scrollHeight")
        driver.set_window_size(BROWSER_WIDTH, min(total_height + 100, 5000))
        time.sleep(0.5)
        inject_nav_css(driver)
        driver.save_screenshot(filepath)
        driver.set_window_size(BROWSER_WIDTH, BROWSER_HEIGHT)
    else:
        driver.save_screenshot(filepath)

    print('  [OK] ' + filename)
    return filepath


def screenshot_element(driver, selector, filename, wait=0.3):
    """Screenshot a specific element."""
    time.sleep(wait)
    filepath = os.path.join(SCREENSHOT_DIR, filename)
    try:
        el = driver.find_element(By.CSS_SELECTOR, selector)
        el.screenshot(filepath)
        print('  [OK] ' + filename + ' (element)')
        return filepath
    except Exception as e:
        print('  [WARN] ' + filename + ': ' + str(e))
        driver.save_screenshot(filepath)
        return filepath


# ============================================================
# PPT HELPERS
# ============================================================
def add_title_slide(prs, title, subtitle=''):
    """Add a title slide with blue gradient-like background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Blue background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    # Decorative shape
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.33), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p2.alignment = PP_ALIGN.LEFT

    return slide


def add_section_slide(prs, section_num, title, desc=''):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)

    # Section number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(2.0), Inches(1.2), Inches(1.2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(10)

    # Title
    txBox = slide.shapes.add_textbox(Inches(2.3), Inches(2.0), Inches(10), Inches(1.2))
    tf2 = txBox.text_frame
    tf2.word_wrap = True
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    if desc:
        txBox3 = slide.shapes.add_textbox(Inches(2.3), Inches(3.5), Inches(10), Inches(1.5))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = desc
        p3.font.size = Pt(18)
        p3.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)

    return slide


def add_content_slide(prs, title, image_path, bullets=None, note=''):
    """Add a content slide with screenshot and optional bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Top bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    bar.line.fill.background()

    # Title in bar
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if bullets and len(bullets) > 0:
        # Layout: left = image (max 8.2" x 6.2"), right = bullets (4.5")
        if image_path and os.path.exists(image_path):
            add_fitted_picture(slide, image_path,
                Inches(0.2), Inches(0.9), Inches(8.2), Inches(6.4))

        txBox2 = slide.shapes.add_textbox(Inches(8.5), Inches(0.9), Inches(4.5), Inches(6.4))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i > 0:
                p2 = tf2.add_paragraph()
            else:
                p2 = tf2.paragraphs[0]
            p2.space_before = Pt(4)
            p2.space_after = Pt(2)

            # Check if it's a header (starts with ##)
            if bullet.startswith('## '):
                p2.text = bullet[3:]
                p2.font.size = Pt(12)
                p2.font.bold = True
                p2.font.color.rgb = DARK_BLUE
                p2.space_before = Pt(10)
            elif bullet.startswith('* '):
                p2.text = bullet[2:]
                p2.font.size = Pt(10)
                p2.font.color.rgb = GRAY
                p2.level = 1
            else:
                # Numbered or icon bullet
                run = p2.add_run()
                run.text = bullet
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(0x33, 0x44, 0x55)
    else:
        # Full-width image (max 12.9" x 6.4")
        if image_path and os.path.exists(image_path):
            add_fitted_picture(slide, image_path,
                Inches(0.2), Inches(0.9), Inches(12.9), Inches(6.4))

    # Note at bottom
    if note:
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12), Inches(0.4))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = note
        p3.font.size = Pt(10)
        p3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        p3.font.italic = True

    return slide


def add_step_slide(prs, step_num, title, image_path, instructions, tip=''):
    """Add a step-by-step instruction slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # Top bar with step number
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    bar.line.fill.background()

    # Step badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.15), Inches(1.2), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_BLUE
    badge.line.fill.background()
    tf_badge = badge.text_frame
    p_badge = tf_badge.paragraphs[0]
    p_badge.text = 'STEP ' + str(step_num)
    p_badge.font.size = Pt(14)
    p_badge.font.bold = True
    p_badge.font.color.rgb = WHITE
    p_badge.alignment = PP_ALIGN.CENTER

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.7), Inches(0.1), Inches(11), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Image on left (max 8.2" x 6.4")
    if image_path and os.path.exists(image_path):
        add_fitted_picture(slide, image_path,
            Inches(0.2), Inches(0.9), Inches(8.2), Inches(6.4))

    # Instructions on right
    txBox2 = slide.shapes.add_textbox(Inches(8.5), Inches(0.9), Inches(4.5), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, inst in enumerate(instructions):
        if i > 0:
            p2 = tf2.add_paragraph()
        else:
            p2 = tf2.paragraphs[0]

        p2.space_before = Pt(3)
        p2.space_after = Pt(2)

        if inst.startswith('>> '):
            # Highlight instruction
            p2.text = inst[3:]
            p2.font.size = Pt(11)
            p2.font.bold = True
            p2.font.color.rgb = ACCENT_BLUE
        elif inst.startswith('! '):
            # Warning
            p2.text = inst[2:]
            p2.font.size = Pt(10)
            p2.font.color.rgb = ACCENT_RED
        else:
            p2.text = inst
            p2.font.size = Pt(10)
            p2.font.color.rgb = RGBColor(0x33, 0x44, 0x55)

    # Tip box
    if tip:
        tip_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(6.4), Inches(4.5), Inches(0.7))
        tip_box.fill.solid()
        tip_box.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
        tip_box.line.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
        tf_tip = tip_box.text_frame
        tf_tip.word_wrap = True
        p_tip = tf_tip.paragraphs[0]
        p_tip.text = 'TIP: ' + tip
        p_tip.font.size = Pt(9)
        p_tip.font.color.rgb = BLUE

    return slide


# ============================================================
# MAIN
# ============================================================
def main():
    print('=== IT 자산 관리 시스템 가이드 PPT 생성 ===')
    print()

    # Setup
    driver = setup_driver()
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    try:
        # ---- Login ----
        print('[1/12] 로그인 페이지...')
        img_login = screenshot(driver, '/login', '01_login.png')
        login(driver)

        # ---- Dashboard ----
        print('[2/12] 대시보드...')
        img_dashboard = screenshot(driver, '/', '02_dashboard.png', wait=2)

        # ---- Inventory (입출고) ----
        print('[3/12] 입출고 관리...')
        img_inv_list = screenshot(driver, '/inventory', '03_inventory_list.png', wait=2)

        # Incoming registration (입고 등록)
        img_incoming = screenshot(driver, '/inventory/incoming', '03_incoming.png', wait=2)

        # Incoming: select server type
        driver.get(BASE_URL + '/inventory/incoming')
        time.sleep(1)
        inject_nav_css(driver)
        try:
            sel = Select(driver.find_element(By.ID, 'incomingType'))
            sel.select_by_value('server')
            time.sleep(1)
            img_incoming_server = screenshot(driver, None, '03_incoming_server.png')
        except:
            img_incoming_server = img_incoming

        # Incoming: vendor ownership
        driver.get(BASE_URL + '/inventory/incoming')
        time.sleep(1)
        inject_nav_css(driver)
        try:
            sel = Select(driver.find_element(By.ID, 'incomingType'))
            sel.select_by_value('server')
            time.sleep(0.5)
            own = Select(driver.find_element(By.ID, 'incomingOwnership'))
            own.select_by_value('vendor')
            time.sleep(0.5)
            img_incoming_vendor = screenshot(driver, None, '03_incoming_vendor.png')
        except:
            img_incoming_vendor = img_incoming

        # Usage registration (사용 등록) - viewport only for top section
        img_inv_new = screenshot(driver, '/inventory/new', '03_inventory_new.png', wait=2)

        # Usage registration: full page (cropped to fit)
        img_inv_new_full = screenshot(driver, '/inventory/new', '03_inventory_new_full.png', wait=2, full_page=True)
        crop_image_if_tall(img_inv_new_full, 1400)

        # Equipment detail
        img_equip_detail = screenshot(driver, '/inventory/equipment/TPC-SV-2U-01', '03_equip_detail.png', wait=2)

        # Inventory list filtered by 입고
        img_inv_incoming_filter = screenshot(driver, '/inventory?status=%EC%9E%85%EA%B3%A0', '03_inv_incoming_filter.png', wait=2)

        # ---- Requests (신청서) ----
        print('[4/12] 신청서...')
        img_requests = screenshot(driver, '/requests', '04_requests.png', wait=2)

        # ---- Assets (자산현황) ----
        print('[5/12] 자산 관리...')
        img_assets = screenshot(driver, '/assets', '05_assets.png', wait=2)
        img_asset_detail = screenshot(driver, '/assets/1107', '05_asset_detail.png', wait=2)
        img_asset_new = screenshot(driver, '/assets/new', '05_asset_new.png', wait=2)
        img_asset_edit = screenshot(driver, '/assets/1107/edit', '05_asset_edit.png', wait=2)
        img_vendor_assets = screenshot(driver, '/assets/vendor', '05_vendor_assets.png', wait=2)

        # ---- Module Inventory (부품현황) ----
        print('[6/12] 부품 현황...')
        img_modules = screenshot(driver, '/module-inventory', '06_modules.png', wait=2)
        img_modules_installed = screenshot(driver, '/module-inventory?tab=installed', '06_modules_installed.png', wait=2)

        # ---- Server Rooms (서버실) ----
        print('[7/12] 서버실...')
        img_rooms = screenshot(driver, '/rooms', '07_rooms.png', wait=2)

        # Find a room with racks
        driver.get(BASE_URL + '/rooms')
        time.sleep(1)
        inject_nav_css(driver)
        try:
            room_link = driver.find_element(By.CSS_SELECTOR, 'a[href*="/rooms/"]')
            room_link.click()
            time.sleep(2)
            img_room_detail = screenshot(driver, None, '07_room_detail.png')
        except:
            img_room_detail = None

        # Racks (server room management view)
        img_racks_rooms = screenshot(driver, '/racks', '07_racks_rooms.png', wait=2)

        # Rack detail - navigate from room detail page
        img_rack_detail = None
        try:
            driver.get(BASE_URL + '/rooms')
            time.sleep(1)
            inject_nav_css(driver)
            room_link = driver.find_element(By.CSS_SELECTOR, 'a[href*="/rooms/"]')
            room_link.click()
            time.sleep(2)
            inject_nav_css(driver)
            rack_link = driver.find_element(By.CSS_SELECTOR, 'a.rack-thumb-link')
            rack_link.click()
            time.sleep(2)
            img_rack_detail = screenshot(driver, None, '07_rack_detail.png')
        except Exception as e:
            print('  [WARN] rack detail: ' + str(e))
            img_rack_detail = None

        # ---- Offices (사무실) ----
        print('[8/12] 사무실/장비실...')
        img_offices = screenshot(driver, '/offices', '08_offices.png', wait=2)
        img_storage = screenshot(driver, '/storage', '08_storage.png', wait=2)

        # ---- Power Panel (전력분전반) ----
        print('[9/12] 전력분전반...')
        img_power = screenshot(driver, '/power-panel', '09_power.png', wait=2)

        # ---- Network Layout ----
        print('[10/12] 네트워크 / IP 관리...')
        img_network = screenshot(driver, '/network-layout', '10_network.png', wait=2)
        img_ip_mgmt = screenshot(driver, '/ip-management', '10_ip_mgmt.png', wait=2)

        # ---- Discovery ----
        img_discovery = screenshot(driver, '/discovery', '10_discovery.png', wait=2)

        # ---- GPU Monitoring ----
        print('[11/12] GPU 모니터링 / 이력 / 백업...')
        img_gpu = screenshot(driver, '/gpu-monitoring', '11_gpu.png', wait=2)
        img_audit = screenshot(driver, '/audit-log', '11_audit.png', wait=2)
        img_backup = screenshot(driver, '/backup', '11_backup.png', wait=2)

        # ---- Admin: Users ----
        print('[12/12] 계정 관리...')
        img_users = screenshot(driver, '/admin/users', '12_users.png', wait=2)

        driver.quit()
        print()
        print('모든 스크린샷 완료. PPT 생성 중...')

        # ============================================================
        # BUILD PPT
        # ============================================================

        # === Cover ===
        add_title_slide(prs,
            'IT 자산 관리 시스템\n사용 가이드',
            '장비 입고부터 반납까지 — 전체 워크플로우 상세 안내\nMLCommons HPC IT Management System')

        # === TOC ===
        slide_toc = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide_toc.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        bar = slide_toc.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK_BLUE
        bar.line.fill.background()

        txBox = slide_toc.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = '목차 (Table of Contents)'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        toc_items = [
            ('1', '로그인', '시스템 접속 및 권한 안내'),
            ('2', '대시보드', '전체 현황 한눈에 보기'),
            ('3', '입출고 관리', '신규: 입고 등록 → 사용 등록 / 기존: 사용 등록 → 반납'),
            ('4', '신청서 관리', '장비 사용 신청서 작성 및 관리'),
            ('5', '자산 현황', '전체 자산 조회 / 등록 / 수정 / 삭제'),
            ('6', '부품 현황', '컴퓨팅 모듈 재고 및 설치 관리'),
            ('7', '서버실 / 랙 관리', '서버실 배치도 및 랙 상세'),
            ('8', '사무실 / 장비실', '사무실 좌석 배치 및 장비실 관리'),
            ('9', '전력 분전반', '전력 노드 및 사용량 관리'),
            ('10', '네트워크 / IP 관리', '네트워크 레이아웃 및 IP 대역 관리'),
            ('11', 'GPU 모니터링 / 이력 / 백업', '실시간 모니터링, 감사 로그, 데이터 백업'),
            ('12', '계정 관리', '사용자 계정 생성 및 권한 설정 (관리자 전용)'),
        ]

        for i, (num, title, desc) in enumerate(toc_items):
            y = 1.1 + i * 0.47
            num_box = slide_toc.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(0.5), Inches(0.38))
            num_box.fill.solid()
            num_box.fill.fore_color.rgb = ACCENT_BLUE
            num_box.line.fill.background()
            ntf = num_box.text_frame
            np = ntf.paragraphs[0]
            np.text = num
            np.font.size = Pt(12)
            np.font.bold = True
            np.font.color.rgb = WHITE
            np.alignment = PP_ALIGN.CENTER

            txB = slide_toc.shapes.add_textbox(Inches(1.5), Inches(y), Inches(3), Inches(0.38))
            ttf = txB.text_frame
            ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tp = ttf.paragraphs[0]
            tp.text = title
            tp.font.size = Pt(14)
            tp.font.bold = True
            tp.font.color.rgb = DARK_BLUE

            txD = slide_toc.shapes.add_textbox(Inches(4.8), Inches(y), Inches(8), Inches(0.38))
            dtf = txD.text_frame
            dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
            dp = dtf.paragraphs[0]
            dp.text = desc
            dp.font.size = Pt(12)
            dp.font.color.rgb = GRAY

        # ============================================================
        # SECTION 1: 로그인
        # ============================================================
        add_section_slide(prs, 1, '로그인', '시스템 접속 및 권한별 기능 안내')

        add_content_slide(prs, '로그인 페이지', img_login, [
            '## 접속 방법',
            '1. 웹 브라우저에서 시스템 URL 접속',
            '2. 사용자명과 비밀번호를 입력',
            '3. "로그인" 버튼 클릭',
            '',
            '## 권한별 기능',
            '- 관리자(admin): 모든 기능 사용 가능',
            '- 유지보수(maintenance): 자산 등록/수정 가능',
            '- 조회(viewer): 읽기 전용',
            '',
            '## 주의사항',
            '- 24시간 후 자동 로그아웃',
            '- 비밀번호 분실 시 관리자에게 문의',
        ])

        # ============================================================
        # SECTION 2: 대시보드
        # ============================================================
        add_section_slide(prs, 2, '대시보드', '전체 현황을 한눈에 파악')

        add_content_slide(prs, '대시보드 — 전체 현황', img_dashboard, [
            '## 주요 정보',
            '- 자산 유형별 현황 (서버, 스토리지, 스위치 등)',
            '- 입출고 상태별 현황',
            '- 최근 등록/변경 이력',
            '',
            '## 네비게이션',
            '- 상단 메뉴바에서 각 기능으로 이동',
            '- 통합 검색: 장비명, IP, 사용자 등 검색',
            '- 관리자만 백업/계정관리 메뉴 표시',
        ])

        # ============================================================
        # SECTION 3: 입출고 관리 (핵심 섹션)
        # ============================================================
        add_section_slide(prs, 3, '입출고 관리',
            '신규 장비: 입고 등록 → 사용 등록 → 반납\n'
            '기존 자산: 사용 등록 → 반납\n\n'
            '업체 장비 / 구매 장비 모두 반드시 "입고 등록"을 먼저 해야 합니다.')

        # ---- 워크플로우 개요 슬라이드 ----
        add_content_slide(prs, '입출고 핵심 워크플로우', img_inv_list, [
            '## 반드시 알아야 할 흐름',
            '',
            '>> [신규 장비가 들어오면]',
            '1) 입고 등록 (먼저!)',
            '   → 자산이 시스템에 생성됨',
            '2) 사용 등록',
            '   → 위치/IP/하드웨어 배정',
            '3) 반납 처리',
            '',
            '>> [기존에 등록된 자산이면]',
            '1) 사용 등록 (바로!)',
            '   → 관리번호 입력 시 자동 로딩',
            '',
            '! 신규 장비를 입고 등록 없이',
            '! 사용 등록하면 안 됩니다!',
            '',
            '## 버튼 위치',
            '- 녹색 "+ 입고 등록": 신규 장비',
            '- 파란색 "+ 사용 등록": 배치/재배치',
        ], note='입출고 목록 상단에 두 버튼이 구분되어 있습니다.')

        # ---- 입고 등록: STEP 1 ----
        add_step_slide(prs, 1, '[입고 등록] 자산유형 및 소유 선택', img_incoming_server, [
            '>> 입출고 → 녹색 "+ 입고 등록" 클릭',
            '',
            '■ 자산유형 선택:',
            '  - 장비: 서버, 스토리지, 스위치 등',
            '  - 부품: CPU, 메모리, 디스크, PSU 등',
            '',
            '■ 소유구분 선택:',
            '  - 자사: 직접 구매한 장비',
            '  - 업체: 외부 업체에서 반입된 장비',
            '',
            '! 모든 신규 장비(구매/업체 불문)는',
            '! 반드시 이 화면에서 입고 등록을 합니다.',
        ], tip='서버 선택 시 노드 수, U 사이즈 등 장비 정보를 추가 입력할 수 있습니다.')

        # ---- 입고 등록: STEP 2 (구매 서버) ----
        add_step_slide(prs, 2, '[입고 등록] 구매 서버 — 기본 정보 입력', img_incoming_server, [
            '>> 소유구분: "자사" 선택',
            '',
            '1. 관리번호 입력 (필수):',
            '   - 관리번호 체계에 맞게 부여',
            '   - 예: TPC-SV-2U-35',
            '   - 추천 번호가 자동 표시됩니다',
            '',
            '2. 자산번호: 자산관리 번호 입력',
            '',
            '3. 제조사 / 모델명 / 시리얼번호 입력',
            '',
            '4. 장비 정보:',
            '   - U 사이즈 (1U, 2U, 4U 등)',
            '   - 입고일, 구매일, 보증만료일',
            '',
            '>> "입고 등록" 버튼 클릭',
        ], tip='입고 등록 완료 후 입출고 목록에서 상태가 "입고"로 표시됩니다.')

        # ---- 입고 등록: STEP 3 (업체 서버) ----
        add_step_slide(prs, 3, '[입고 등록] 업체 장비 — 업체 선택 및 등록', img_incoming_vendor, [
            '>> 소유구분: "업체" 선택',
            '',
            '1. 업체 선택:',
            '   - 기존 등록된 업체 선택',
            '   - 또는 "+ 신규 업체 등록"으로 새 업체 추가',
            '',
            '2. 관리번호 자동 생성:',
            '   - 업체명 기반으로 자동 부여',
            '   - 예: 코코링크-001, 글루시스-003',
            '',
            '3. 나머지 정보 입력 (구매 서버와 동일):',
            '   - 모델명, 시리얼번호, U 사이즈',
            '   - 반납예정일 (업체 장비의 경우)',
            '',
            '>> "입고 등록" 버튼 클릭',
        ], tip='업체 장비는 반납예정일을 입력해두면 관리가 편리합니다.')

        # ---- 입고 완료 확인 ----
        add_content_slide(prs, '입고 등록 완료 확인', img_inv_incoming_filter, [
            '## 입고 등록 후',
            '',
            '>> 입출고 목록에서 "입고" 탭 클릭',
            '',
            '1. 방금 등록한 장비가',
            '   "입고" 상태로 표시됩니다.',
            '',
            '2. 이 시점에서:',
            '   - 자산현황에 자산이 자동 생성됨',
            '   - 아직 위치/IP 미배정 상태',
            '',
            '## 다음 단계',
            '>> 이제 "사용 등록"을 진행합니다.',
            '   (위치, IP, 하드웨어, 접속정보 배정)',
        ], note='입고 상태의 장비는 아직 서버실에 배치되지 않은 상태입니다.')

        # ---- 사용 등록 시작 ----
        add_step_slide(prs, 4, '[사용 등록] 기본 정보 — 기존 자산 로딩', img_inv_new_full, [
            '>> 입출고 → 파란색 "+ 사용 등록" 클릭',
            '',
            '■ 기존 자산 (입고 완료된 장비):',
            '  - 관리번호 입력 → 자동완성 목록',
            '  - 선택 시 기존 정보 자동 로딩:',
            '    모델명, IP, 접속정보, 하드웨어 등',
            '',
            '■ 기존 자산 (이미 사용중인 장비):',
            '  - 위치 변경, 정보 업데이트 시 사용',
            '  - 관리번호 입력만으로 모든 정보 로딩',
            '',
            '! 신규 장비는 반드시 입고 등록을 먼저 하고',
            '! 이 화면에서 사용 등록을 진행하세요.',
        ], tip='관리번호 또는 자산번호 어느 것으로든 자동완성이 됩니다.')

        # ---- 사용 등록: 하드웨어 ----
        add_step_slide(prs, 5, '[사용 등록] 하드웨어 정보 입력', img_inv_new_full, [
            '>> 하드웨어 섹션에서 부품 추가',
            '',
            '1. "+ CPU", "+ 메모리" 등 버튼으로 추가',
            '',
            '2. 각 행 입력:',
            '   - 소유: 자사 / 업체',
            '   - 유형: CPU, 메모리, 디스크, NIC 등',
            '   - 부품 코드: 드롭다운 또는 직접 입력',
            '   - 수량: 장착 수량',
            '',
            '3. PSU (전원장치):',
            '   - "+ PSU" 버튼으로 추가',
            '   - Active / Standby 역할 선택',
            '',
            '4. 기존 자산 로딩 시:',
            '   - 하드웨어가 자동으로 채워집니다.',
            '   - 필요 시 수정/추가/삭제 가능',
        ], tip='업체 소유 부품은 "업체" 선택 시 모델명 직접 입력으로 전환됩니다.')

        # ---- 사용 등록: IP / 접속정보 / 위치 ----
        add_step_slide(prs, 6, '[사용 등록] IP / 접속정보 / 위치', img_inv_new_full, [
            '>> IP, 접속정보, 위치를 순서대로 입력',
            '',
            '■ IP 주소:',
            '  - "+ IP" 버튼으로 행 추가',
            '  - 용도: Management / BMC / IB 등',
            '  - 인터페이스 유형 및 속도 선택',
            '',
            '■ 접속 정보:',
            '  - "+ 계정" 버튼으로 추가',
            '  - 유형: Root / Admin / 기타',
            '',
            '■ 위치:',
            '  - 서버실 선택 → 랙 선택',
            '  - 랙 미리보기에서 빈 위치 클릭',
            '  - 기존 자산의 현재 위치는 빈 칸 표시',
            '',
            '>> 모든 입력 후 "등록" 버튼 클릭',
        ], tip='랙 미리보기에서 자기 자신의 위치는 비어있게 표시됩니다.')

        # ---- 장비 이력 상세 ----
        add_content_slide(prs, '장비 이력 상세', img_equip_detail, [
            '## 확인 가능한 정보',
            '- 기본 정보 (관리번호, 모델명 등)',
            '- 하드웨어 구성 (CPU, 메모리 등)',
            '- IP 주소 목록',
            '- 접속 정보 (계정)',
            '- 입출고 이력 (시간순)',
            '',
            '## 접근 방법',
            '- 입출고 목록에서 관리번호 클릭',
            '- 또는 자산 상세에서 "상세 보기"',
            '',
            '## 상태 흐름 요약',
            '>> 입고 → 사용중 → 반납완료',
            '(필요 시 재입고/재사용 가능)',
        ])

        # ============================================================
        # SECTION 4: 신청서
        # ============================================================
        add_section_slide(prs, 4, '신청서 관리', '장비 사용 신청서 작성 및 승인 관리')

        add_content_slide(prs, '신청서 목록 및 작성', img_requests, [
            '## 기능',
            '- 장비 사용 신청서 작성',
            '- 신청 상태 관리 (대기/승인/반려)',
            '- 신청 이력 조회',
            '',
            '## 사용 방법',
            '1. "신청서" 메뉴 클릭',
            '2. "새 신청서" 버튼으로 작성',
            '3. 필요한 장비 정보 입력',
            '4. 관리자가 승인/반려 처리',
        ])

        # ============================================================
        # SECTION 5: 자산 현황
        # ============================================================
        add_section_slide(prs, 5, '자산 현황',
            '전체 자산 조회, 등록, 수정, 삭제\n'
            '업체 장비와 자사 장비를 통합 관리합니다.')

        add_content_slide(prs, '자산 현황 — 목록', img_assets, [
            '## 필터 기능',
            '- 유형: 서버, 스토리지, 스위치 등',
            '- 소유: 자사 / 업체',
            '- 상태: 활성 / 유지보수 / 비활성 등',
            '- 서버실별 필터',
            '- 텍스트 검색',
            '',
            '## 정렬',
            '- 각 컬럼 헤더 클릭으로 정렬',
            '- 관리번호 클릭 시 상세 페이지 이동',
            '',
            '## 관리 버튼',
            '- 수정: 자산 정보 편집',
            '- 삭제: 자산 삭제 (확인 필요)',
        ])

        add_content_slide(prs, '자산 상세 페이지', img_asset_detail, [
            '## 표시 정보',
            '- 기본 정보 (자산번호, 관리번호, 유형 등)',
            '- 위치 정보 (서버실, 랙, U 위치)',
            '- IP 주소 목록',
            '- 접속 정보',
            '- 운영 정보 (사용자, 용도)',
            '',
            '## 컴퓨팅 모듈',
            '- CPU, 메모리, 디스크, NIC, GPU, PSU',
            '- 소유 구분 (자사/업체명 표시)',
            '- 설치 유형 (온보드/애드온)',
            '- "+ 모듈 추가" 버튼으로 직접 추가',
            '',
            '## 입출고 이력',
            '- 해당 장비의 전체 사용/반납 이력',
        ])

        add_content_slide(prs, '업체 장비 현황', img_vendor_assets, [
            '## 업체별 장비 보기',
            '- 자산현황 → "업체 장비" 버튼 클릭',
            '- 업체별로 그룹화하여 표시',
            '- 각 업체의 장비 수량 및 상세 확인',
            '',
            '## 업체 장비 등록 방법',
            '1. 입출고 → 사용등록',
            '2. 소유구분: "업체" 선택',
            '3. 관리번호 자동 생성',
            '   (업체명-001, 002, ...)',
            '',
            '## 활용',
            '- 업체별 장비 현황 파악',
            '- 반납 시 업체별 정리 가능',
        ])

        add_content_slide(prs, '자산 등록 / 수정', img_asset_edit, [
            '## 자산 등록',
            '- 자산현황 목록에서 직접 등록 가능',
            '- 입출고 사용등록과는 별도 경로',
            '',
            '## 자산 수정',
            '- 목록에서 "수정" 또는 상세에서 "수정"',
            '- 기본정보, 위치, IP, 접속정보 편집',
            '',
            '## 랙 미리보기',
            '- 랙 선택 시 슬롯 레벨 미리보기 표시',
            '- 3홀/U (하/중/상) 단위로 표시',
            '- 빈 위치 클릭으로 자동 배치',
            '- 현재 편집 중인 자산은 비어있게 표시',
            '',
            '## 선반 설정',
            '- "선반 포함" 체크 → 선반 U 크기 입력',
            '- 장비 아래에 선반 공간 자동 추가',
        ])

        # ============================================================
        # SECTION 6: 부품 현황
        # ============================================================
        add_section_slide(prs, 6, '부품 현황', '컴퓨팅 모듈 재고 관리 및 설치 현황')

        add_content_slide(prs, '부품 현황 — 재고 / 설치', img_modules, [
            '## 탭 구분',
            '- 재고: 부품 재고 목록 (미설치)',
            '- 설치됨: 자산에 설치된 모듈 현황',
            '',
            '## 재고 관리',
            '- 부품 코드별 재고 수량 관리',
            '- 입고 / 출고 / 설치 / 제거 이력',
            '- 부품 유형: CPU, 메모리, 디스크,',
            '  NIC, RAID, GPU, PSU',
            '',
            '## 모듈 추가',
            '- 재고 탭에서 새 부품 등록',
            '- 자산 상세에서 직접 모듈 추가',
            '- 입출고 사용등록 시 하드웨어 섹션에서',
            '  자동 연동',
        ])

        # ============================================================
        # SECTION 7: 서버실 / 랙 관리
        # ============================================================
        add_section_slide(prs, 7, '서버실 / 랙 관리',
            '서버실 배치도, 랙 상세, 자산 배치 현황')

        add_content_slide(prs, '서버실 목록', img_rooms, [
            '## 서버실 관리',
            '- 서버실 목록 및 상세 정보',
            '- 위치 유형: 서버실 / 사무실 / 장비실',
            '- 서버실별 랙 배치 현황',
            '',
            '## 랙 관리',
            '- 랙별 자산 배치 시각화',
            '- U 단위 사용 현황',
            '- 자산 수 / 사용량 표시',
            '',
            '## 서버실 상세',
            '- 랙 배치도 (행/열 그리드)',
            '- 각 랙 클릭 시 상세 정보',
            '- 전면/후면 보기',
        ])

        if img_rack_detail:
            add_content_slide(prs, '랙 상세', img_rack_detail, [
                '## 랙 상세 보기',
                '- U 단위 자산 배치 시각화',
                '- 전면/후면 전환',
                '- 블레이드 서버 좌/우 구분',
                '',
                '## 자산 정보',
                '- 각 자산 클릭 시 상세 페이지 이동',
                '- 컴퓨팅 모듈 탭: 장착 부품 확인',
                '- IP 주소 탭: 네트워크 정보',
                '',
                '## 빈 슬롯',
                '- 빈 슬롯 확인 가능',
                '- 자산 등록/수정 시 랙 미리보기 연동',
            ])

        # ============================================================
        # SECTION 8: 사무실 / 장비실
        # ============================================================
        add_section_slide(prs, 8, '사무실 / 장비실', '좌석 배치 및 장비실 관리')

        add_content_slide(prs, '사무실 관리', img_offices, [
            '## 사무실',
            '- 사무실 목록 및 좌석 배치',
            '- 인원 배치 현황',
            '',
            '## 장비실',
            '- 서버실 외 장비 보관 공간',
            '- 부품 보관 위치 관리',
        ])

        add_content_slide(prs, '장비실 관리', img_storage, [
            '## 장비실 현황',
            '- 장비실별 보관 장비 목록',
            '- 부품 보관 위치 관리',
            '',
            '## 활용',
            '- 미사용 장비 보관 관리',
            '- 여분 부품 위치 추적',
        ])

        # ============================================================
        # SECTION 9: 전력 분전반
        # ============================================================
        add_section_slide(prs, 9, '전력 분전반', '전력 노드 및 사용량 관리')

        add_content_slide(prs, '전력 분전반', img_power, [
            '## 기능',
            '- 전력 분전반 구성도',
            '- 전력 노드별 연결 장비 관리',
            '- 전력 사용량 모니터링',
            '',
            '## 관리 방법',
            '1. 분전반 노드 추가',
            '2. 각 노드에 장비 연결',
            '3. 전력 사용 현황 모니터링',
        ])

        # ============================================================
        # SECTION 10: 네트워크 / IP 관리
        # ============================================================
        add_section_slide(prs, 10, '네트워크 / IP 관리',
            '네트워크 레이아웃, IP 대역 관리, 자산 매칭')

        add_content_slide(prs, '네트워크 레이아웃', img_network, [
            '## 기능',
            '- 서버실별 네트워크 구성도',
            '- 스위치 ↔ 자산 연결 관계',
            '- 포트별 연결 상태 시각화',
            '',
            '## 관리 방법',
            '1. 스위치 자산을 네트워크에 등록',
            '2. 각 포트에 자산 연결',
            '3. 연결 구성도 확인',
        ])

        add_content_slide(prs, 'IP 관리', img_ip_mgmt, [
            '## 기능',
            '- IP 서브넷 대역 관리',
            '- IP별 할당 상태 확인',
            '- 가용 IP 자동 탐지',
            '',
            '## 관리 방법',
            '1. 서브넷 대역 등록 (예: 10.100.0.0/24)',
            '2. IP 할당 현황 확인',
            '3. 자산에 IP 할당 시 자동 동기화',
        ])

        add_content_slide(prs, '자산 매칭 (Discovery)', img_discovery, [
            '## 기능',
            '- SSH 기반 서버 하드웨어 자동 탐지',
            '- 탐지 결과 ↔ 등록 자산 비교',
            '- 하드웨어 변경 감지',
            '',
            '## 사용 방법',
            '1. IP 대역 또는 개별 IP 입력',
            '2. 접속 정보 입력 (또는 자산 정보 활용)',
            '3. "탐색" 실행',
            '4. 결과와 기존 자산 매칭/비교',
        ])

        # ============================================================
        # SECTION 11: GPU 모니터링 / 이력 / 백업
        # ============================================================
        add_section_slide(prs, 11, 'GPU 모니터링 / 이력 / 백업',
            '실시간 GPU 상태, 감사 로그, 데이터 백업')

        add_content_slide(prs, 'GPU 모니터링', img_gpu, [
            '## 기능',
            '- GPU 장착 서버 실시간 모니터링',
            '- GPU 온도, 사용률, 메모리 현황',
            '- nvidia-smi 정보 수집',
            '',
            '## 활용',
            '- GPU 서버 상태 한눈에 파악',
            '- 이상 징후 조기 발견',
        ])

        add_content_slide(prs, '이력 관리 (감사 로그)', img_audit, [
            '## 기능',
            '- 모든 자산 변경 이력 기록',
            '- 사용자별 작업 이력 조회',
            '- 날짜/유형별 필터링',
            '',
            '## 기록 항목',
            '- 자산 생성/수정/삭제',
            '- 입출고 등록/반납',
            '- 로그인/로그아웃',
            '- 부품 이동/변경',
        ])

        add_content_slide(prs, '백업 관리 (관리자 전용)', img_backup, [
            '## 기능',
            '- 데이터베이스 수동 백업',
            '- 자동 백업 (7일 주기)',
            '- 백업 파일 다운로드 및 복원',
            '',
            '## 자동 백업',
            '- 서버 시작 시 7일 주기로 자동 실행',
            '- 1년 경과한 백업 자동 삭제',
            '',
            '! 관리자 권한 필요',
        ])

        # ============================================================
        # SECTION 12: 계정 관리
        # ============================================================
        add_section_slide(prs, 12, '계정 관리 (관리자 전용)',
            '사용자 계정 생성, 수정, 삭제 및 권한 설정')

        add_content_slide(prs, '계정 관리', img_users, [
            '## 접근 방법',
            '- 관리자로 로그인 후 "계정관리" 메뉴',
            '- 관리자 외에는 메뉴 미표시 + 접근 차단',
            '',
            '## 계정 생성',
            '1. "새 계정 추가" 폼에서 입력:',
            '   - 사용자명 (로그인 ID)',
            '   - 비밀번호',
            '   - 표시 이름',
            '   - 권한: 관리자 / 유지보수 / 조회',
            '2. "계정 생성" 버튼 클릭',
            '',
            '## 권한 설명',
            '- admin: 모든 기능 (계정/백업 포함)',
            '- maintenance: 자산 등록/수정/삭제',
            '- viewer: 읽기 전용',
            '',
            '! 현재 로그인된 계정은 삭제 불가',
        ])

        # ============================================================
        # END: 마무리
        # ============================================================
        end_slide = add_title_slide(prs,
            '감사합니다',
            'IT 자산 관리 시스템 사용 가이드\n\n'
            '문의사항은 시스템 관리자에게 연락해주세요.')

        # Save
        prs.save(OUTPUT_PPT)
        print()
        print('PPT 생성 완료: ' + OUTPUT_PPT)
        print('총 슬라이드 수: ' + str(len(prs.slides)))

    except Exception as e:
        import traceback
        traceback.print_exc()
        try:
            driver.quit()
        except:
            pass
        sys.exit(1)


if __name__ == '__main__':
    main()
